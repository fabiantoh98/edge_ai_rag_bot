from abc import ABC, abstractmethod
import csv
from PIL import Image
import whisper
from pptx import Presentation
from PyPDF2 import PdfReader
from paddleocr import PaddleOCR
import json

class BaseExtractor(ABC):
    @abstractmethod
    def extract(self, file_path: str) -> list[str]:
        """Extracts text content from file and returns list of text chunks or lines."""
        pass

class AudioVideoExtractor(BaseExtractor):
    """_summary_
    * Do not pass in video or audio files that does not contain speech
    Perform transcribing for audio and video files.
    In this project we will not be using each frame in the video.
    
    Args:
        BaseExtractor (_type_): _description_
    """
    def __init__(self, model_name="base"):
        self.model = whisper.load_model(model_name)

    def extract(self, file_path: str) -> list[str]:
        result = self.model.transcribe(file_path)
        return result["text"]

class CSVExtractor(BaseExtractor):
    def extract(self, file_path: str) -> list[str]:
        texts = []
        with open(file_path, newline='', encoding='utf-8') as csvfile:
            reader = csv.reader(csvfile)
            for row in reader:
                # Join all columns in the row into a single string separated by comma
                texts.append(",".join(str(cell) for cell in row if cell))
        return texts
    
# May not need since we are using multi modal embedding and llm
class ImageExtractor(BaseExtractor):
    """_summary_
    Args:
        BaseExtractor (_type_): _description_
    """
    def __init__(self, model_name="base"):
        self.ocr = PaddleOCR(
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
            use_textline_orientation=False)

    def extract(self, file_path: str) -> list[str]:
        result = self.ocr.predict(input=file_path)
        return ", ".join(result[0]["rec_texts"])
    
class PDFExtractor(BaseExtractor):
    def extract(self, file_path: str) -> list[str]:
        reader = PdfReader(file_path)
        texts=[]
        for page in reader.pages:
            text = page.extract_text()
            texts.append(text)
        return texts
    
# todo: Account for the actual slide number based on the raw pptx
    
class PPTXExtractor(BaseExtractor):
    """
    Extracts unique slide texts from a PPTX file, preserving order.
    Each slide's text is concatenated, duplicates are removed, and the result is a list.
    The index in the list corresponds to the slide number (for metadata).
    """
    def extract(self, file_path: str) -> list[str]:
        prs = Presentation(file_path)
        seen = set()
        unique_texts = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if " " in shape.text.strip():
                        slide_text.append(shape.text.strip())
            joined = " ".join(set(slide_text))
            if joined and joined not in seen:
                seen.add(joined)
                unique_texts.append(joined)
        return unique_texts

def load_qa_from_json(filename):
    with open(filename, "r", encoding="utf-8") as f:
        qa_list = json.load(f)
    questions = [item["question"] for item in qa_list]
    ground_truth = [item["ground_truth"] for item in qa_list]
    return questions, ground_truth